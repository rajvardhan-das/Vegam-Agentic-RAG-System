from pathlib import Path
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

def load_pdf(path: Path) -> str:
    reader = PdfReader(path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def load_docx(path: Path) -> str:
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])


def load_pptx(path: Path) -> str:
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides_text.append(shape.text)
    return "\n".join(slides_text)





def load_excel(path: Path) -> str:
    wb = load_workbook(path, data_only=True)
    output = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        output.append(f"Sheet: {sheet_name}")

        headers = None
        for row in ws.iter_rows(values_only=True):
            if headers is None:
                headers = row
                continue
            row_text = " ; ".join(
                f"{headers[i]}: {row[i]}" for i in range(len(headers))
            )
            output.append(row_text)


    return "\n".join(output)



def load_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def load_document(file_path: str) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return load_pdf(path)
    elif suffix == ".docx":
        return load_docx(path)
    elif suffix == ".pptx":
        return load_pptx(path)
    elif suffix in [".xls", ".xlsx"]:
        return load_excel(path)
    elif suffix == ".txt":
        return load_txt(path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")
